"""
PowerPoint document processor.

This module handles extraction of text, images, tables, and metadata from PowerPoint files.
"""

import logging
from typing import Dict, List, Any, Optional
from pathlib import Path
import base64
import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

logger = logging.getLogger(__name__)


class PowerPointProcessor:
    """PowerPoint document processor using python-pptx."""
    
    def __init__(self):
        self.supported_formats = ['.ppt', '.pptx']
    
    async def process(self, file_path: str) -> Dict[str, Any]:
        """
        Process PowerPoint file and extract all content.
        
        Returns:
            Dict containing extracted content, metadata, and structure
        """
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")
            
            # Open PowerPoint presentation
            prs = Presentation(str(file_path))
            
            # Extract basic metadata
            metadata = self._extract_metadata(prs)
            
            # Extract content from all slides
            slides = []
            all_text = ""
            images = []
            tables = []
            charts = []
            
            for slide_num, slide in enumerate(prs.slides):
                # Extract slide content
                slide_data = self._extract_slide_content(slide, slide_num)
                
                slides.append(slide_data)
                all_text += slide_data["text"] + "\n\n"
                images.extend(slide_data["images"])
                tables.extend(slide_data["tables"])
                charts.extend(slide_data["charts"])
            
            # Analyze presentation structure
            structure = self._analyze_structure(slides)
            
            # Extract brand elements
            brand_elements = self._extract_brand_elements(slides, all_text)
            
            return {
                "content_type": "powerpoint",
                "metadata": metadata,
                "slide_count": len(slides),
                "word_count": len(all_text.split()),
                "text": all_text,
                "slides": slides,
                "images": images,
                "tables": tables,
                "charts": charts,
                "structure": structure,
                "brand_elements": brand_elements,
            }
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {e}")
            raise
    
    def _extract_metadata(self, prs: Presentation) -> Dict[str, Any]:
        """Extract PowerPoint metadata."""
        try:
            core_props = prs.core_properties
            return {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "comments": core_props.comments or "",
                "category": core_props.category or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "last_modified_by": core_props.last_modified_by or "",
                "slide_count": len(prs.slides),
            }
        except Exception as e:
            logger.warning(f"Error extracting metadata: {e}")
            return {"slide_count": len(prs.slides)}
    
    def _extract_slide_content(self, slide, slide_num: int) -> Dict[str, Any]:
        """Extract content from a single slide."""
        slide_data = {
            "slide_number": slide_num + 1,
            "text": "",
            "title": "",
            "content": [],
            "images": [],
            "tables": [],
            "charts": [],
            "layout": "",
            "notes": "",
        }
        
        try:
            # Get slide layout name
            slide_data["layout"] = slide.slide_layout.name
            
            # Extract notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    slide_data["notes"] = notes_slide.notes_text_frame.text
            
            # Process all shapes on the slide
            text_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Identify title (usually the first text shape or largest font)
                    if not slide_data["title"] and len(text) < 100:
                        slide_data["title"] = text
                    
                    text_content.append(text)
                    slide_data["content"].append({
                        "type": "text",
                        "content": text,
                        "shape_type": str(shape.shape_type),
                    })
                
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Extract image
                    image_data = self._extract_image(shape, slide_num)
                    if image_data:
                        slide_data["images"].append(image_data)
                
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    # Extract table
                    table_data = self._extract_table(shape, slide_num)
                    if table_data:
                        slide_data["tables"].append(table_data)
                
                elif shape.shape_type in [MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT]:
                    # Extract chart
                    chart_data = self._extract_chart(shape, slide_num)
                    if chart_data:
                        slide_data["charts"].append(chart_data)
            
            slide_data["text"] = "\n".join(text_content)
            slide_data["word_count"] = len(slide_data["text"].split())
            
        except Exception as e:
            logger.warning(f"Error extracting content from slide {slide_num}: {e}")
        
        return slide_data
    
    def _extract_image(self, shape, slide_num: int) -> Optional[Dict[str, Any]]:
        """Extract image from a shape."""
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Convert to PIL Image for processing
            pil_image = Image.open(io.BytesIO(image_bytes))
            
            # Convert to base64 for storage
            buffer = io.BytesIO()
            pil_image.save(buffer, format="PNG")
            img_base64 = base64.b64encode(buffer.getvalue()).decode()
            
            return {
                "slide_number": slide_num + 1,
                "width": pil_image.width,
                "height": pil_image.height,
                "format": pil_image.format or "PNG",
                "data": img_base64,
                "size_bytes": len(image_bytes),
                "filename": getattr(image, 'filename', f"slide_{slide_num}_image"),
            }
            
        except Exception as e:
            logger.warning(f"Error extracting image from slide {slide_num}: {e}")
            return None
    
    def _extract_table(self, shape, slide_num: int) -> Optional[Dict[str, Any]]:
        """Extract table from a shape."""
        try:
            table = shape.table
            table_data = []
            
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    row_data.append(cell_text)
                table_data.append(row_data)
            
            return {
                "slide_number": slide_num + 1,
                "rows": len(table_data),
                "columns": len(table_data[0]) if table_data else 0,
                "data": table_data,
            }
            
        except Exception as e:
            logger.warning(f"Error extracting table from slide {slide_num}: {e}")
            return None
    
    def _extract_chart(self, shape, slide_num: int) -> Optional[Dict[str, Any]]:
        """Extract chart information from a shape."""
        try:
            # Basic chart information extraction
            chart_data = {
                "slide_number": slide_num + 1,
                "type": "chart",
                "shape_type": str(shape.shape_type),
            }
            
            # Try to extract chart data if available
            if hasattr(shape, 'chart'):
                chart = shape.chart
                chart_data.update({
                    "chart_type": str(chart.chart_type) if hasattr(chart, 'chart_type') else "unknown",
                    "has_title": hasattr(chart, 'chart_title') and chart.chart_title is not None,
                })
                
                # Extract chart title if available
                if hasattr(chart, 'chart_title') and chart.chart_title:
                    chart_data["title"] = chart.chart_title.text_frame.text
            
            return chart_data
            
        except Exception as e:
            logger.warning(f"Error extracting chart from slide {slide_num}: {e}")
            return {"slide_number": slide_num + 1, "type": "chart", "error": str(e)}
    
    def _analyze_structure(self, slides: List[Dict]) -> Dict[str, Any]:
        """Analyze presentation structure."""
        structure = {
            "sections": [],
            "slide_types": {},
            "flow": [],
        }
        
        try:
            current_section = None
            
            for slide in slides:
                # Categorize slide types
                layout = slide.get("layout", "unknown")
                if layout not in structure["slide_types"]:
                    structure["slide_types"][layout] = 0
                structure["slide_types"][layout] += 1
                
                # Detect section breaks (slides with short titles, specific layouts)
                title = slide.get("title", "")
                if title and len(title.split()) <= 5 and layout in ["Title Slide", "Section Header"]:
                    # Start new section
                    if current_section:
                        structure["sections"].append(current_section)
                    
                    current_section = {
                        "title": title,
                        "start_slide": slide["slide_number"],
                        "slides": [slide["slide_number"]],
                    }
                elif current_section:
                    current_section["slides"].append(slide["slide_number"])
                
                # Track presentation flow
                structure["flow"].append({
                    "slide": slide["slide_number"],
                    "title": title,
                    "layout": layout,
                    "content_length": len(slide.get("text", "")),
                })
            
            # Add last section
            if current_section:
                structure["sections"].append(current_section)
            
        except Exception as e:
            logger.warning(f"Error analyzing presentation structure: {e}")
        
        return structure
    
    def _extract_brand_elements(self, slides: List[Dict], all_text: str) -> Dict[str, Any]:
        """Extract brand-specific elements from the presentation."""
        brand_elements = {
            "colors": [],
            "fonts": [],
            "logos": [],
            "brand_mentions": [],
            "design_patterns": [],
            "slide_templates": [],
        }
        
        try:
            # Analyze slide layouts for design patterns
            layout_counts = {}
            for slide in slides:
                layout = slide.get("layout", "unknown")
                layout_counts[layout] = layout_counts.get(layout, 0) + 1
            
            brand_elements["slide_templates"] = [
                {"layout": layout, "count": count}
                for layout, count in layout_counts.items()
            ]
            
            # Look for brand-related keywords
            import re
            
            brand_keywords = [
                "brand", "logo", "color", "font", "typography", "voice", "tone",
                "guidelines", "identity", "palette", "style", "design", "mission",
                "vision", "values", "positioning"
            ]
            
            for keyword in brand_keywords:
                mentions = []
                for slide in slides:
                    slide_text = slide.get("text", "")
                    if keyword.lower() in slide_text.lower():
                        # Find specific mentions
                        lines = slide_text.split('\n')
                        for line in lines:
                            if keyword.lower() in line.lower():
                                mentions.append({
                                    "keyword": keyword,
                                    "slide": slide["slide_number"],
                                    "context": line.strip(),
                                })
                
                if mentions:
                    brand_elements["brand_mentions"].extend(mentions)
            
            # Identify potential logos (small images, specific positions)
            for slide in slides:
                for img in slide.get("images", []):
                    # Heuristic: small images might be logos
                    if img["width"] < 300 and img["height"] < 300:
                        brand_elements["logos"].append({
                            "slide": img["slide_number"],
                            "size": f"{img['width']}x{img['height']}",
                            "format": img["format"],
                        })
            
            # Look for color mentions in text
            hex_colors = re.findall(r'#[0-9A-Fa-f]{6}', all_text)
            rgb_colors = re.findall(r'rgb\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*\)', all_text, re.IGNORECASE)
            
            brand_elements["colors"].extend(hex_colors)
            brand_elements["colors"].extend(rgb_colors)
            
        except Exception as e:
            logger.warning(f"Error extracting brand elements: {e}")
        
        return brand_elements

"""
Brand Playbook Intelligence Engine
Complete vector-based semantic search and passage retrieval system
"""

import logging
import os
import uuid
import asyncio
import json
import time
from typing import List, Dict, Any, Optional, Tuple
import chromadb
from chromadb.config import Settings
import openai
import tiktoken
from sentence_transformers import SentenceTransformer
import PyPDF2
from pptx import Presentation
import io
import re
from dataclasses import dataclass
from supabase import create_client, Client
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize clients
supabase: Client = create_client(
    os.getenv("SUPABASE_URL"),
    os.getenv("SUPABASE_SERVICE_KEY")
)

# Initialize ChromaDB
chroma_client = chromadb.PersistentClient(
    path="./chroma_db",
    settings=Settings(
        anonymized_telemetry=False,
        allow_reset=True
    )
)

# Initialize embedding model
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

# Initialize tokenizer for text chunking
tokenizer = tiktoken.get_encoding("cl100k_base")

@dataclass
class DocumentChunk:
    """Represents a chunk of text from a document"""
    id: str
    text: str
    document_id: str
    document_name: str
    page_number: int
    chunk_index: int
    token_count: int
    metadata: Dict[str, Any]

@dataclass
class SearchResult:
    """Represents a search result with source attribution"""
    passage: str
    document_name: str
    page_number: int
    relevance_score: float
    document_id: str
    chunk_id: str

@dataclass
class IntelligentResponse:
    """Complete AI response with sources and confidence"""
    answer: str
    confidence: float
    sources: List[SearchResult]
    query: str
    processing_time: float

class DocumentProcessor:
    """Handles document processing and text extraction"""
    
    @staticmethod
    async def extract_pdf_text_with_pages(content: bytes) -> List[Tuple[str, int]]:
        """Extract text from PDF with page numbers"""
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            pages_text = []
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():  # Only add non-empty pages
                    pages_text.append((text.strip(), page_num))
            
            return pages_text
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            return [("Error extracting PDF text", 1)]
    
    @staticmethod
    async def extract_ppt_text_with_slides(content: bytes) -> List[Tuple[str, int]]:
        """Extract text from PowerPoint with slide numbers"""
        try:
            ppt_file = io.BytesIO(content)
            presentation = Presentation(ppt_file)
            slides_text = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():  # Only add non-empty slides
                    slides_text.append((slide_text.strip(), slide_num))
            
            return slides_text
        except Exception as e:
            logger.error(f"PPT extraction error: {e}")
            return [("Error extracting PowerPoint text", 1)]

class TextChunker:
    """Handles intelligent text chunking with overlap"""
    
    @staticmethod
    def chunk_text(text: str, max_tokens: int = 200, overlap_tokens: int = 30) -> List[str]:
        """Split text into smaller, more focused overlapping chunks"""
        # Clean and preprocess text
        text = re.sub(r'\s+', ' ', text.strip())

        # Try to split by sentences first for better semantic chunks
        sentences = re.split(r'[.!?]+', text)
        chunks = []
        current_chunk = ""

        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue

            # Check if adding this sentence would exceed token limit
            test_chunk = current_chunk + " " + sentence if current_chunk else sentence
            tokens = tokenizer.encode(test_chunk)

            if len(tokens) <= max_tokens:
                current_chunk = test_chunk
            else:
                # Save current chunk if it has content
                if current_chunk:
                    chunks.append(current_chunk)

                # Start new chunk with current sentence
                current_chunk = sentence

                # If single sentence is too long, split by tokens
                if len(tokenizer.encode(current_chunk)) > max_tokens:
                    token_chunks = TextChunker._split_by_tokens(current_chunk, max_tokens, overlap_tokens)
                    chunks.extend(token_chunks)
                    current_chunk = ""

        # Add final chunk
        if current_chunk:
            chunks.append(current_chunk)

        return chunks

    @staticmethod
    def _split_by_tokens(text: str, max_tokens: int, overlap_tokens: int) -> List[str]:
        """Fallback: split by tokens when sentences are too long"""
        tokens = tokenizer.encode(text)
        chunks = []

        start = 0
        while start < len(tokens):
            end = min(start + max_tokens, len(tokens))
            chunk_tokens = tokens[start:end]
            chunk_text = tokenizer.decode(chunk_tokens)
            chunks.append(chunk_text)

            if end >= len(tokens):
                break

            start = end - overlap_tokens

        return chunks
    
    @staticmethod
    def create_document_chunks(
        pages_text: List[Tuple[str, int]], 
        document_id: str, 
        document_name: str
    ) -> List[DocumentChunk]:
        """Create document chunks with metadata"""
        chunks = []
        chunk_index = 0
        
        for text, page_number in pages_text:
            # Split page text into chunks
            page_chunks = TextChunker.chunk_text(text)
            
            for chunk_text in page_chunks:
                chunk_id = f"{document_id}_chunk_{chunk_index}"
                token_count = len(tokenizer.encode(chunk_text))
                
                chunk = DocumentChunk(
                    id=chunk_id,
                    text=chunk_text,
                    document_id=document_id,
                    document_name=document_name,
                    page_number=page_number,
                    chunk_index=chunk_index,
                    token_count=token_count,
                    metadata={
                        "document_type": "pdf" if document_name.endswith('.pdf') else "ppt",
                        "created_at": str(uuid.uuid4())
                    }
                )
                chunks.append(chunk)
                chunk_index += 1
        
        return chunks

class VectorDatabase:
    """Manages ChromaDB operations"""
    
    def __init__(self, collection_name: str = "brand_playbooks"):
        self.collection_name = collection_name
        self.collection = None
        self._initialize_collection()
    
    def _initialize_collection(self):
        """Initialize or get existing collection"""
        try:
            self.collection = chroma_client.get_or_create_collection(
                name=self.collection_name,
                metadata={"description": "Brand playbook document chunks"}
            )
            logger.info(f"Initialized collection: {self.collection_name}")
        except Exception as e:
            logger.error(f"Error initializing collection: {e}")
            raise
    
    async def add_document_chunks(self, chunks: List[DocumentChunk], user_id: str):
        """Add document chunks to vector database"""
        try:
            # Prepare data for ChromaDB
            ids = [chunk.id for chunk in chunks]
            documents = [chunk.text for chunk in chunks]
            metadatas = []
            
            for chunk in chunks:
                metadata = {
                    "document_id": chunk.document_id,
                    "document_name": chunk.document_name,
                    "page_number": chunk.page_number,
                    "chunk_index": chunk.chunk_index,
                    "token_count": chunk.token_count,
                    "user_id": user_id,
                    **chunk.metadata
                }
                metadatas.append(metadata)
            
            # Generate embeddings
            embeddings = embedding_model.encode(documents).tolist()
            
            # Add to ChromaDB
            self.collection.add(
                ids=ids,
                documents=documents,
                metadatas=metadatas,
                embeddings=embeddings
            )
            
            logger.info(f"Added {len(chunks)} chunks to vector database")
            return True
            
        except Exception as e:
            logger.error(f"Error adding chunks to vector database: {e}")
            return False
    
    async def search_similar_passages(
        self, 
        query: str, 
        user_id: str, 
        n_results: int = 5,
        document_ids: Optional[List[str]] = None
    ) -> List[SearchResult]:
        """Search for similar passages using vector similarity"""
        try:
            # Generate query embedding
            query_embedding = embedding_model.encode([query]).tolist()[0]
            
            # Prepare where clause for user filtering
            where_clause = {"user_id": user_id}
            if document_ids:
                where_clause["document_id"] = {"$in": document_ids}
            
            # Search in ChromaDB
            results = self.collection.query(
                query_embeddings=[query_embedding],
                n_results=n_results,
                where=where_clause,
                include=["documents", "metadatas", "distances"]
            )
            
            # Convert to SearchResult objects
            search_results = []
            if results['documents'] and results['documents'][0]:
                for i, (doc, metadata, distance) in enumerate(zip(
                    results['documents'][0],
                    results['metadatas'][0],
                    results['distances'][0]
                )):
                    # Convert distance to similarity score (0-1)
                    # ChromaDB uses cosine distance, so smaller distance = higher similarity
                    # Normalize distance to similarity score
                    if distance is not None:
                        # For cosine distance, range is typically 0-2, convert to 0-1 similarity
                        relevance_score = max(0, min(1, 1 - (distance / 2)))
                    else:
                        relevance_score = 0.5  # Default if distance is None

                    search_result = SearchResult(
                        passage=doc,
                        document_name=metadata['document_name'],
                        page_number=metadata['page_number'],
                        relevance_score=relevance_score,
                        document_id=metadata['document_id'],
                        chunk_id=metadata.get('chunk_id', f"chunk_{i}")
                    )
                    search_results.append(search_result)

            # Sort by relevance score (highest first)
            search_results.sort(key=lambda x: x.relevance_score, reverse=True)
            
            return search_results
            
        except Exception as e:
            logger.error(f"Error searching vector database: {e}")
            return []
    
    async def get_user_document_count(self, user_id: str) -> int:
        """Get count of documents for a user"""
        try:
            results = self.collection.get(
                where={"user_id": user_id},
                include=["metadatas"]
            )
            
            # Count unique documents
            document_ids = set()
            if results['metadatas']:
                for metadata in results['metadatas']:
                    document_ids.add(metadata['document_id'])
            
            return len(document_ids)
            
        except Exception as e:
            logger.error(f"Error getting document count: {e}")
            return 0
    
    async def delete_document_chunks(self, document_id: str, user_id: str):
        """Delete all chunks for a document"""
        try:
            self.collection.delete(
                where={
                    "document_id": document_id,
                    "user_id": user_id
                }
            )
            logger.info(f"Deleted chunks for document: {document_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error deleting document chunks: {e}")
            return False

class AIResponseGenerator:
    """Generates intelligent responses using retrieved passages"""

    @staticmethod
    async def generate_response(
        query: str,
        search_results: List[SearchResult],
        api_key: str,
        provider: str = "openai"
    ) -> IntelligentResponse:
        """Generate AI response with source attribution"""
        import time
        start_time = time.time()

        try:
            # Prepare context from search results
            context_passages = []
            for i, result in enumerate(search_results[:3], 1):  # Use top 3 results
                context_passages.append(
                    f"Source {i} (from {result.document_name}, page {result.page_number}):\n{result.passage}\n"
                )

            context = "\n".join(context_passages)

            # Create prompt for AI
            system_prompt = """You are an expert brand consultant analyzing brand playbook documents. Your job is to provide accurate, specific answers based ONLY on the provided brand playbook content.

IMPORTANT RULES:
1. ONLY use information explicitly stated in the provided sources
2. Be specific and cite exact details (colors, fonts, values, etc.)
3. If the exact information isn't in the sources, clearly state that
4. Reference which source contains the information
5. Provide actionable, detailed answers when possible
6. If multiple sources contain related info, synthesize them clearly"""

            user_prompt = f"""BRAND PLAYBOOK CONTENT:
{context}

QUESTION: {query}

INSTRUCTIONS: Analyze the provided brand playbook content and answer the question with specific details. If the information exists in the sources, provide exact details (like specific colors, fonts, values, guidelines, etc.). If not found, clearly state what information is missing."""

            # Generate response based on provider
            if provider == "openai":
                response_text = await AIResponseGenerator._call_openai(
                    system_prompt, user_prompt, api_key
                )
            else:  # claude
                response_text = await AIResponseGenerator._call_claude(
                    system_prompt, user_prompt, api_key
                )

            # Calculate confidence based on relevance scores
            if search_results:
                avg_relevance = sum(r.relevance_score for r in search_results) / len(search_results)
                confidence = min(0.95, avg_relevance * 1.2)  # Cap at 95%
            else:
                confidence = 0.1

            processing_time = time.time() - start_time

            return IntelligentResponse(
                answer=response_text,
                confidence=confidence,
                sources=search_results,
                query=query,
                processing_time=processing_time
            )

        except Exception as e:
            logger.error(f"Error generating AI response: {e}")
            processing_time = time.time() - start_time

            return IntelligentResponse(
                answer=f"I apologize, but I encountered an error while processing your question: {str(e)}",
                confidence=0.0,
                sources=search_results,
                query=query,
                processing_time=processing_time
            )

    @staticmethod
    async def _call_openai(system_prompt: str, user_prompt: str, api_key: str) -> str:
        """Call OpenAI API"""
        try:
            client = openai.AsyncOpenAI(api_key=api_key)

            response = await client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_tokens=500,
                temperature=0.3
            )

            return response.choices[0].message.content

        except Exception as e:
            logger.error(f"OpenAI API error: {e}")
            return f"Error calling OpenAI API: {str(e)}"

    @staticmethod
    async def _call_claude(system_prompt: str, user_prompt: str, api_key: str) -> str:
        """Call Claude API"""
        try:
            import httpx

            async with httpx.AsyncClient() as client:
                response = await client.post(
                    "https://api.anthropic.com/v1/messages",
                    headers={
                        "x-api-key": api_key,
                        "Content-Type": "application/json",
                        "anthropic-version": "2023-06-01"
                    },
                    json={
                        "model": "claude-3-sonnet-20240229",
                        "max_tokens": 500,
                        "messages": [
                            {"role": "user", "content": f"{system_prompt}\n\n{user_prompt}"}
                        ]
                    },
                    timeout=30.0
                )

                if response.status_code == 200:
                    return response.json()["content"][0]["text"]
                else:
                    return f"Claude API error: {response.status_code}"

        except Exception as e:
            logger.error(f"Claude API error: {e}")
            return f"Error calling Claude API: {str(e)}"

    @staticmethod
    async def generate_enhanced_response(
        query: str,
        search_results: List[SearchResult],
        api_key: str,
        provider: str = "openai",
        feedback: str = ""
    ) -> IntelligentResponse:
        """Generate enhanced response with improved prompting based on feedback"""
        start_time = time.time()

        try:
            # Create enhanced context with better formatting
            context_parts = []
            for i, result in enumerate(search_results, 1):
                context_parts.append(
                    f"SOURCE {i} (Relevance: {result.relevance_score:.2f}):\n"
                    f"Document: {result.document_name}\n"
                    f"Page: {result.page_number}\n"
                    f"Content: {result.passage}\n"
                    f"---"
                )

            context = "\n".join(context_parts)

            # Enhanced system prompt targeting 90% quality threshold
            system_prompt = f"""You are an expert brand consultant with deep expertise in brand strategy, visual identity, and brand guidelines. Your job is to provide EXCEPTIONAL answers that meet a 90% quality threshold based on brand playbook documents.

CRITICAL QUALITY REQUIREMENTS (Must achieve 90%+):

1. ACCURACY (90%+): Every statement must be factually correct and directly address the question
2. COMPLETENESS (90%+): Provide comprehensive, actionable information covering all aspects
3. RELEVANCE (90%+): Every detail must be directly relevant to the question asked
4. SPECIFICITY (90%+): Include exact details - hex codes, font names, measurements, specifications
5. SOURCE USAGE (90%+): Properly cite and utilize all relevant sources

MANDATORY INSTRUCTIONS:
- ONLY use information explicitly stated in the provided sources
- Be extremely specific: exact colors (#HEX codes), precise font names, measurements, percentages
- For colors: provide exact color names, hex codes, RGB values, usage contexts
- For fonts: provide exact typeface names, weights, sizes, usage guidelines
- For brand elements: provide specific implementation details, dimensions, spacing
- Always cite which source contains each piece of information
- If information is missing, clearly state what specific details are not available
- Structure answers for maximum clarity and actionability

FEEDBACK TO ADDRESS: {feedback if feedback else "Achieve 90% quality across all criteria"}

QUALITY THRESHOLD: This answer must score 90%+ or it will be rejected. Be comprehensive, specific, and accurate."""

            user_prompt = f"""BRAND PLAYBOOK SOURCES:
{context}

QUESTION: {query}

INSTRUCTIONS:
1. Analyze all provided sources carefully
2. Extract specific, actionable information that directly answers the question
3. Include exact details like color codes, font names, measurements, etc.
4. Cite which source(s) contain each piece of information
5. If the question asks for specific elements (colors, fonts, etc.) but they're not in the sources, explicitly state what's missing
6. Provide a comprehensive answer that a brand manager could immediately use

ANSWER:"""

            # Call AI with enhanced prompts
            if provider == "openai":
                ai_response = await AIResponseGenerator._call_openai_enhanced(system_prompt, user_prompt, api_key)
            else:
                ai_response = await AIResponseGenerator._call_claude(system_prompt, user_prompt, api_key)

            # Calculate confidence based on source quality and relevance
            avg_relevance = sum(r.relevance_score for r in search_results) / len(search_results)
            source_count_factor = min(len(search_results) / 5, 1.0)  # Normalize to max 5 sources
            confidence = (avg_relevance * 0.7) + (source_count_factor * 0.3)

            processing_time = time.time() - start_time

            return IntelligentResponse(
                answer=ai_response,
                confidence=confidence,
                sources=search_results,
                query=query,
                processing_time=processing_time
            )

        except Exception as e:
            logger.error(f"Error generating enhanced AI response: {e}")
            processing_time = time.time() - start_time

            return IntelligentResponse(
                answer=f"I apologize, but I encountered an error while processing your question: {str(e)}",
                confidence=0.0,
                sources=search_results,
                query=query,
                processing_time=processing_time
            )

    @staticmethod
    async def _call_openai_enhanced(system_prompt: str, user_prompt: str, api_key: str) -> str:
        """Call OpenAI API with enhanced settings for better responses"""
        try:
            client = openai.AsyncOpenAI(api_key=api_key)

            response = await client.chat.completions.create(
                model="gpt-4",  # Use GPT-4 for better quality
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_tokens=800,  # Allow longer responses
                temperature=0.1,  # Lower temperature for more focused responses
                top_p=0.9
            )

            return response.choices[0].message.content

        except Exception as e:
            logger.error(f"Enhanced OpenAI API error: {e}")
            return f"Error calling OpenAI API: {str(e)}"

class PlaybookIntelligence:
    """Main intelligence engine orchestrator"""

    def __init__(self):
        self.vector_db = VectorDatabase()
        self.document_processor = DocumentProcessor()
        self.text_chunker = TextChunker()
        self.ai_generator = AIResponseGenerator()

    async def process_document(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        user_id: str
    ) -> bool:
        """Complete document processing pipeline"""
        try:
            logger.info(f"Processing document: {filename}")

            # Extract text with page/slide numbers
            if filename.lower().endswith('.pdf'):
                pages_text = await self.document_processor.extract_pdf_text_with_pages(content)
            else:  # PowerPoint
                pages_text = await self.document_processor.extract_ppt_text_with_slides(content)

            # Create chunks
            chunks = self.text_chunker.create_document_chunks(
                pages_text, document_id, filename
            )

            # Store in vector database
            success = await self.vector_db.add_document_chunks(chunks, user_id)

            if success:
                logger.info(f"Successfully processed {len(chunks)} chunks from {filename}")
                return True
            else:
                logger.error(f"Failed to store chunks for {filename}")
                return False

        except Exception as e:
            logger.error(f"Error processing document {filename}: {e}")
            return False

    async def answer_question_enhanced(
        self,
        query: str,
        user_id: str,
        api_key: str,
        provider: str = "openai",
        document_ids: Optional[List[str]] = None
    ) -> IntelligentResponse:
        """Enhanced answer generation with LLM-based quality evaluation and improvement"""
        try:
            logger.info(f"Processing enhanced query: {query}")

            # Stage 1: Initial search and answer generation
            initial_response = await self._generate_initial_answer(query, user_id, api_key, provider, document_ids)

            # Stage 2: LLM evaluation of the answer quality
            evaluation = await self._evaluate_answer_quality(query, initial_response, api_key, provider)

            # Stage 3: Iteratively improve until we reach 90% quality threshold
            current_response = initial_response
            current_evaluation = evaluation
            improvement_attempts = 0
            max_attempts = 3

            while current_evaluation["quality_score"] < 0.9 and improvement_attempts < max_attempts:
                logger.info(f"Answer quality ({current_evaluation['quality_score']:.2f}) below 90% threshold, attempting improvement #{improvement_attempts + 1}")

                improved_response = await self._improve_answer(query, user_id, api_key, provider, document_ids, current_evaluation["feedback"])
                improved_evaluation = await self._evaluate_answer_quality(query, improved_response, api_key, provider)

                if improved_evaluation["quality_score"] > current_evaluation["quality_score"]:
                    logger.info(f"Improvement successful: {current_evaluation['quality_score']:.2f} → {improved_evaluation['quality_score']:.2f}")
                    current_response = improved_response
                    current_evaluation = improved_evaluation
                else:
                    logger.info(f"Improvement attempt #{improvement_attempts + 1} did not increase quality")

                improvement_attempts += 1

            # Final quality check - only return if we meet the 90% threshold
            if current_evaluation["quality_score"] >= 0.9:
                logger.info(f"✅ High-quality answer achieved: {current_evaluation['quality_score']:.2f}")
                return current_response
            else:
                logger.warning(f"❌ Could not achieve 90% quality after {max_attempts} attempts. Best score: {current_evaluation['quality_score']:.2f}")

                # Return a response indicating insufficient quality
                return IntelligentResponse(
                    answer=f"I apologize, but I cannot provide a sufficiently accurate answer to your question based on the available brand playbook content. The information may not be detailed enough in the uploaded documents, or the question may require information not present in the current playbooks.\n\nTo get a better answer, you might:\n1. Upload additional brand documentation\n2. Rephrase your question to be more specific\n3. Ask about topics that are more thoroughly covered in the existing playbooks\n\nCurrent quality assessment: {current_evaluation['quality_score']:.1%} (requires 90% for delivery)",
                    confidence=current_evaluation["quality_score"],
                    sources=current_response.sources,
                    query=query,
                    processing_time=current_response.processing_time
                )

        except Exception as e:
            logger.error(f"Error in enhanced query processing: {e}")
            # Fallback to original method
            return await self.answer_question(query, user_id, api_key, provider, document_ids)

    async def answer_question(
        self,
        query: str,
        user_id: str,
        api_key: str,
        provider: str = "openai",
        document_ids: Optional[List[str]] = None
    ) -> IntelligentResponse:
        """Answer question using intelligent search and AI"""
        try:
            logger.info(f"Processing query: {query}")

            # Search for relevant passages with more results for better coverage
            search_results = await self.vector_db.search_similar_passages(
                query, user_id, n_results=10, document_ids=document_ids
            )

            # Filter out very low relevance results (below 0.3 threshold)
            search_results = [r for r in search_results if r.relevance_score > 0.3]

            if not search_results:
                return IntelligentResponse(
                    answer="I couldn't find any relevant information in your uploaded brand playbooks to answer this question. Please make sure you have uploaded your brand documents first.",
                    confidence=0.0,
                    sources=[],
                    query=query,
                    processing_time=0.0
                )

            # Generate AI response
            response = await self.ai_generator.generate_response(
                query, search_results, api_key, provider
            )

            return response

        except Exception as e:
            logger.error(f"Error answering question: {e}")
            return IntelligentResponse(
                answer=f"I encountered an error while processing your question: {str(e)}",
                confidence=0.0,
                sources=[],
                query=query,
                processing_time=0.0
            )

    async def _generate_initial_answer(self, query: str, user_id: str, api_key: str, provider: str, document_ids: Optional[List[str]] = None) -> IntelligentResponse:
        """Generate initial answer using standard RAG approach"""
        return await self.answer_question(query, user_id, api_key, provider, document_ids)

    async def _evaluate_answer_quality(self, query: str, response: IntelligentResponse, api_key: str, provider: str) -> Dict[str, Any]:
        """Use LLM to evaluate the quality of the generated answer"""
        try:
            evaluation_prompt = f"""You are an expert evaluator of AI-generated answers about brand playbooks and marketing materials. You have VERY HIGH STANDARDS and only accept exceptional answers.

ORIGINAL QUESTION: {query}

AI GENERATED ANSWER: {response.answer}

SOURCES USED: {len(response.sources)} sources with confidence {response.confidence:.2f}

EVALUATION CRITERIA (Each must score 0.9+ for overall 90% quality):

1. ACCURACY (0.0-1.0): Is the answer factually correct and directly addresses the question?
   - 1.0: Perfect accuracy, directly answers the question
   - 0.9: Highly accurate with minor gaps
   - 0.8: Mostly accurate but some inaccuracies
   - <0.8: Significant accuracy issues

2. COMPLETENESS (0.0-1.0): Does it provide comprehensive, actionable information?
   - 1.0: Complete answer with all necessary details
   - 0.9: Very comprehensive, minor details missing
   - 0.8: Good coverage but missing some important elements
   - <0.8: Incomplete or superficial

3. RELEVANCE (0.0-1.0): Is every piece of information directly relevant to the question?
   - 1.0: Every detail is perfectly relevant
   - 0.9: Highly relevant with minimal tangential content
   - 0.8: Mostly relevant but some off-topic content
   - <0.8: Contains irrelevant information

4. SPECIFICITY (0.0-1.0): Does it include exact details (hex codes, font names, measurements, etc.)?
   - 1.0: Provides exact specifications when available
   - 0.9: Very specific with minor gaps in precision
   - 0.8: Good specificity but could be more precise
   - <0.8: Vague or lacks specific details

5. SOURCE_USAGE (0.0-1.0): Does it properly cite and utilize the available sources?
   - 1.0: Perfect source attribution and utilization
   - 0.9: Excellent source usage with minor gaps
   - 0.8: Good source usage but could be better
   - <0.8: Poor source utilization

QUALITY THRESHOLD: Only answers scoring 0.9+ overall should be delivered to users.

Respond in JSON format:
{{
    "quality_score": 0.0-1.0,
    "accuracy": 0.0-1.0,
    "completeness": 0.0-1.0,
    "relevance": 0.0-1.0,
    "specificity": 0.0-1.0,
    "source_usage": 0.0-1.0,
    "feedback": "Specific, actionable feedback on what must be improved to reach 90% quality",
    "missing_elements": ["specific", "elements", "needed", "for", "90%", "quality"],
    "meets_threshold": true/false
}}"""

            if provider == "openai":
                client = openai.OpenAI(api_key=api_key)
                eval_response = client.chat.completions.create(
                    model="gpt-4",
                    messages=[
                        {"role": "system", "content": "You are an expert evaluator. Respond only with valid JSON."},
                        {"role": "user", "content": evaluation_prompt}
                    ],
                    temperature=0.1
                )

                evaluation_text = eval_response.choices[0].message.content
                evaluation = json.loads(evaluation_text)

                logger.info(f"Answer evaluation: quality_score={evaluation['quality_score']:.2f}")
                return evaluation

        except Exception as e:
            logger.error(f"Error evaluating answer quality: {e}")
            # Return default evaluation if LLM evaluation fails
            return {
                "quality_score": 0.5,
                "accuracy": 0.5,
                "completeness": 0.5,
                "relevance": 0.5,
                "specificity": 0.5,
                "source_usage": 0.5,
                "feedback": "Could not evaluate answer quality",
                "missing_elements": []
            }

    async def _improve_answer(self, query: str, user_id: str, api_key: str, provider: str, document_ids: Optional[List[str]], feedback: str) -> IntelligentResponse:
        """Aggressively improve the answer to reach 90% quality threshold"""
        try:
            logger.info(f"Attempting to improve answer for 90% quality. Feedback: {feedback}")

            # Multi-strategy search for comprehensive coverage
            improved_search_results = []

            # Strategy 1: Expand search with related terms and synonyms
            expanded_query = await self._expand_query_with_synonyms(query, api_key, provider)
            expanded_results = await self.vector_db.search_similar_passages(
                expanded_query, user_id, n_results=20, document_ids=document_ids
            )
            improved_search_results.extend(expanded_results)

            # Strategy 2: Targeted searches for specific brand elements
            brand_elements = [
                "color", "colour", "palette", "hex", "RGB", "CMYK",
                "font", "typography", "typeface", "text", "heading",
                "logo", "brand mark", "symbol", "icon",
                "guideline", "standard", "specification", "rule",
                "value", "mission", "vision", "purpose", "essence",
                "voice", "tone", "personality", "character",
                "positioning", "strategy", "message", "tagline"
            ]

            for element in brand_elements:
                if element.lower() in query.lower() or element.lower() in feedback.lower():
                    element_results = await self.vector_db.search_similar_passages(
                        f"{query} {element}", user_id, n_results=15, document_ids=document_ids
                    )
                    improved_search_results.extend(element_results)

            # Strategy 3: Broad context search for comprehensive understanding
            context_queries = [
                f"brand {query}",
                f"visual {query}",
                f"design {query}",
                f"marketing {query}",
                f"identity {query}"
            ]

            for context_query in context_queries:
                context_results = await self.vector_db.search_similar_passages(
                    context_query, user_id, n_results=10, document_ids=document_ids
                )
                improved_search_results.extend(context_results)

            # Remove duplicates and apply higher quality filtering
            seen_chunks = set()
            unique_results = []
            for result in improved_search_results:
                if result.chunk_id not in seen_chunks and result.relevance_score > 0.2:  # Lower threshold for more coverage
                    unique_results.append(result)
                    seen_chunks.add(result.chunk_id)

            # Sort by relevance and take top results for comprehensive coverage
            unique_results.sort(key=lambda x: x.relevance_score, reverse=True)
            final_results = unique_results[:20]  # Use more results for 90% quality

            if not final_results:
                # If no improved results, return original answer
                return await self.answer_question(query, user_id, api_key, provider, document_ids)

            # Generate improved response with enhanced prompt targeting 90% quality
            improved_response = await self.ai_generator.generate_enhanced_response(
                query, final_results, api_key, provider, f"CRITICAL: Must achieve 90% quality. {feedback}"
            )

            return improved_response

        except Exception as e:
            logger.error(f"Error improving answer: {e}")
            # Fallback to original method
            return await self.answer_question(query, user_id, api_key, provider, document_ids)

    async def _expand_query_with_synonyms(self, query: str, api_key: str, provider: str) -> str:
        """Use LLM to expand query with relevant synonyms and related terms"""
        try:
            expansion_prompt = f"""Given this brand playbook question: "{query}"

Generate 3-5 related terms or synonyms that might help find relevant information in brand documents.
Focus on brand-specific terminology.

Examples:
- "brand colors" → "color palette, brand palette, primary colors, secondary colors, color scheme"
- "typography" → "fonts, typeface, text style, font family, heading fonts"
- "brand voice" → "tone of voice, communication style, brand personality, messaging"

Question: {query}
Related terms (comma-separated):"""

            if provider == "openai":
                client = openai.OpenAI(api_key=api_key)
                response = client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "system", "content": "You are a brand expert. Provide only the related terms, comma-separated."},
                        {"role": "user", "content": expansion_prompt}
                    ],
                    temperature=0.3,
                    max_tokens=100
                )

                related_terms = response.choices[0].message.content.strip()
                expanded_query = f"{query} {related_terms}"
                logger.info(f"Expanded query: {expanded_query}")
                return expanded_query

        except Exception as e:
            logger.error(f"Error expanding query: {e}")

        return query  # Return original if expansion fails

    async def get_user_stats(self, user_id: str) -> Dict[str, Any]:
        """Get user statistics"""
        try:
            document_count = await self.vector_db.get_user_document_count(user_id)

            return {
                "documents_processed": document_count,
                "vector_database_status": "active",
                "embedding_model": "all-MiniLM-L6-v2"
            }

        except Exception as e:
            logger.error(f"Error getting user stats: {e}")
            return {
                "documents_processed": 0,
                "vector_database_status": "error",
                "embedding_model": "unknown"
            }

# Initialize global intelligence engine
intelligence_engine = PlaybookIntelligence()

"""
PlaybookWiz Production API with Supabase Integration
"""

import logging
import os
import uuid
import asyncio
from typing import Optional, List
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Depends, Header, Response
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel
from dotenv import load_dotenv
from supabase import create_client, Client
from cryptography.fernet import Fernet
import httpx
import PyPDF2
from pptx import Presentation
import io
import base64

# Load environment variables
load_dotenv()

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Supabase client
supabase: Client = create_client(
    os.getenv("SUPABASE_URL"),
    os.getenv("SUPABASE_SERVICE_KEY")
)

# Encryption for API keys
encryption_key = os.getenv("ENCRYPTION_KEY", "playbookwiz-32char-encryption-key")
# Generate a proper Fernet key
fernet_key = base64.urlsafe_b64encode(encryption_key.encode().ljust(32, b'0')[:32])
cipher_suite = Fernet(fernet_key)

# Pydantic models
class ChatRequest(BaseModel):
    message: str
    session_id: Optional[str] = None
    document_ids: List[str] = []

class IdeationRequest(BaseModel):
    prompt: str
    document_ids: List[str] = []
    personas: List[str] = []

class APIKeyRequest(BaseModel):
    provider: str  # 'openai' or 'claude'
    api_key: str

# FastAPI app
app = FastAPI(
    title="PlaybookWiz Production API",
    description="AI-powered brand playbook analysis and ideation",
    version="1.0.0"
)

# CORS
allowed_origins = ["http://localhost:9000", "http://localhost:3000", "http://127.0.0.1:9000", "http://127.0.0.1:3000"]
app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["*"],
)

# Security
security = HTTPBearer()

async def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)):
    """Verify JWT token and get current user"""
    try:
        logger.info(f"Authenticating user with token: {credentials.credentials[:20]}...")
        # Verify token with Supabase
        user = supabase.auth.get_user(credentials.credentials)
        if not user.user:
            logger.error("No user found in token")
            raise HTTPException(status_code=401, detail="Invalid token")
        logger.info(f"User authenticated: {user.user.id}")
        return user.user
    except Exception as e:
        logger.error(f"Auth error: {e}")
        raise HTTPException(status_code=401, detail="Authentication failed")

# Helper functions
async def extract_pdf_text(content: bytes) -> str:
    """Extract text from PDF"""
    try:
        pdf_file = io.BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        return f"Error extracting PDF text: {str(e)}"

async def extract_ppt_text(content: bytes) -> str:
    """Extract text from PowerPoint"""
    try:
        ppt_file = io.BytesIO(content)
        presentation = Presentation(ppt_file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PPT extraction error: {e}")
        return f"Error extracting PowerPoint text: {str(e)}"

async def get_user_api_keys(user_id: str) -> dict:
    """Get encrypted API keys for user"""
    try:
        result = supabase.table("user_api_keys").select("*").eq("user_id", user_id).execute()
        keys = {}
        for row in result.data:
            try:
                decrypted_key = cipher_suite.decrypt(row["encrypted_key"].encode()).decode()
                keys[row["provider"]] = decrypted_key
            except Exception as e:
                logger.error(f"Decryption error for {row['provider']}: {e}")
        return keys
    except Exception as e:
        logger.error(f"Error getting API keys: {e}")
        return {}

async def call_openai_api(message: str, context: str, api_key: str) -> str:
    """Call OpenAI API"""
    try:
        async with httpx.AsyncClient() as client:
            response = await client.post(
                "https://api.openai.com/v1/chat/completions",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={
                    "model": "gpt-3.5-turbo",
                    "messages": [
                        {"role": "system", "content": "You are a helpful brand expert assistant."},
                        {"role": "user", "content": f"Context: {context}\n\nQuestion: {message}"}
                    ],
                    "max_tokens": 500,
                    "temperature": 0.7
                },
                timeout=30.0
            )
            if response.status_code == 200:
                return response.json()["choices"][0]["message"]["content"]
            else:
                return f"OpenAI API error: {response.status_code}"
    except Exception as e:
        return f"Error calling OpenAI: {str(e)}"

async def call_claude_api(message: str, context: str, api_key: str) -> str:
    """Call Claude API"""
    try:
        async with httpx.AsyncClient() as client:
            response = await client.post(
                "https://api.anthropic.com/v1/messages",
                headers={
                    "x-api-key": api_key,
                    "Content-Type": "application/json",
                    "anthropic-version": "2023-06-01"
                },
                json={
                    "model": "claude-3-sonnet-20240229",
                    "max_tokens": 500,
                    "messages": [{"role": "user", "content": f"Context: {context}\n\nQuestion: {message}"}]
                },
                timeout=30.0
            )
            if response.status_code == 200:
                return response.json()["content"][0]["text"]
            else:
                return f"Claude API error: {response.status_code}"
    except Exception as e:
        return f"Error calling Claude: {str(e)}"

# API Routes
@app.get("/health")
async def health_check():
    return {"status": "healthy", "service": "PlaybookWiz Production API"}

@app.options("/api/v1/auth/api-keys")
async def options_api_keys():
    return Response(
        content="",
        headers={
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "POST, OPTIONS",
            "Access-Control-Allow-Headers": "*",
        }
    )

@app.options("/api/v1/documents/upload")
async def options_upload():
    return Response(
        content="",
        headers={
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "POST, OPTIONS",
            "Access-Control-Allow-Headers": "*",
        }
    )

@app.options("/api/v1/chat/message")
async def options_chat():
    return Response(
        content="",
        headers={
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "POST, OPTIONS",
            "Access-Control-Allow-Headers": "*",
        }
    )

@app.post("/api/v1/auth/api-keys")
async def save_api_key(request: APIKeyRequest, user = Depends(get_current_user)):
    """Save encrypted API key for user"""
    try:
        logger.info(f"Saving API key for user {user.id}, provider: {request.provider}")
        encrypted_key = cipher_suite.encrypt(request.api_key.encode()).decode()

        # Upsert API key (update if exists, insert if not)
        result = supabase.table("user_api_keys").upsert({
            "user_id": user.id,
            "provider": request.provider,
            "encrypted_key": encrypted_key
        }, on_conflict="user_id,provider").execute()

        logger.info(f"API key saved successfully: {result}")
        return {"message": f"{request.provider} API key saved successfully"}
    except Exception as e:
        logger.error(f"Error saving API key: {e}")
        raise HTTPException(status_code=500, detail="Failed to save API key")

@app.post("/api/v1/documents/upload")
async def upload_document(
    file: UploadFile = File(...),
    user = Depends(get_current_user)
):
    """Upload and process document"""
    try:
        # Validate file
        allowed_types = ["application/pdf", "application/vnd.ms-powerpoint", 
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation"]
        allowed_extensions = ['.pdf', '.ppt', '.pptx']
        
        file_extension = '.' + file.filename.split('.')[-1].lower() if '.' in file.filename else None
        
        if file.content_type not in allowed_types and file_extension not in allowed_extensions:
            raise HTTPException(status_code=400, detail="Unsupported file type")
        
        # Read and process file
        content = await file.read()
        
        # Extract text based on file type
        logger.info(f"Extracting text from {file.filename}, type: {file.content_type}, extension: {file_extension}")
        if file.content_type == "application/pdf" or file_extension == '.pdf':
            extracted_text = await extract_pdf_text(content)
            logger.info(f"PDF text extracted: {len(extracted_text)} characters")
        elif file_extension in ['.ppt', '.pptx']:
            extracted_text = await extract_ppt_text(content)
            logger.info(f"PPT text extracted: {len(extracted_text)} characters")
        else:
            extracted_text = "Text extraction not supported for this file type"
            logger.warning(f"Unsupported file type: {file.content_type}, extension: {file_extension}")
        
        # Store in database
        document_data = {
            "user_id": user.id,
            "filename": file.filename,
            "content_type": file.content_type,
            "size_bytes": len(content),
            "extracted_text": extracted_text,
            "status": "processed"
        }
        
        result = supabase.table("documents").insert(document_data).execute()
        document_id = result.data[0]["id"]
        
        return {
            "document_id": document_id,
            "filename": file.filename,
            "status": "processed",
            "message": "Document uploaded and processed successfully"
        }
        
    except Exception as e:
        logger.error(f"Upload error: {e}")
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")

@app.get("/api/v1/documents")
async def get_documents(user = Depends(get_current_user)):
    """Get user's documents"""
    try:
        result = supabase.table("documents").select("id, filename, status, created_at").eq("user_id", user.id).execute()
        return {"documents": result.data}
    except Exception as e:
        logger.error(f"Error getting documents: {e}")
        raise HTTPException(status_code=500, detail="Failed to get documents")

@app.post("/api/v1/chat/message")
async def chat_message(request: ChatRequest, user = Depends(get_current_user)):
    """Process chat message with AI"""
    try:
        # Get user's API keys
        api_keys = await get_user_api_keys(user.id)
        
        if not api_keys:
            raise HTTPException(status_code=400, detail="No API keys configured. Please add your OpenAI or Claude API key.")
        
        # Get document context
        context = ""
        if request.document_ids:
            docs_result = supabase.table("documents").select("filename, extracted_text").in_("id", request.document_ids).eq("user_id", user.id).execute()
            for doc in docs_result.data:
                context += f"\n\nDocument: {doc['filename']}\n{doc['extracted_text']}"
        
        # Call AI API
        if "openai" in api_keys:
            response = await call_openai_api(request.message, context, api_keys["openai"])
        elif "claude" in api_keys:
            response = await call_claude_api(request.message, context, api_keys["claude"])
        else:
            raise HTTPException(status_code=400, detail="No valid API keys found")
        
        # Save chat message
        if request.session_id:
            session_id = request.session_id
        else:
            # Create new session
            session_result = supabase.table("chat_sessions").insert({
                "user_id": user.id,
                "title": request.message[:50] + "..." if len(request.message) > 50 else request.message
            }).execute()
            session_id = session_result.data[0]["id"]
        
        # Save messages
        supabase.table("chat_messages").insert([
            {"session_id": session_id, "role": "user", "content": request.message, "document_ids": request.document_ids},
            {"session_id": session_id, "role": "assistant", "content": response, "document_ids": request.document_ids}
        ]).execute()
        
        return {
            "response": response,
            "session_id": session_id,
            "documents_used": len(request.document_ids)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Chat error: {e}")
        raise HTTPException(status_code=500, detail=f"Chat failed: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("production_main:app", host="0.0.0.0", port=8000, reload=True)
